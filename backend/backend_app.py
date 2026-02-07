import os
import uuid
import time
import threading
import json
import datetime
from functools import wraps

import bcrypt
import fitz  # PyMuPDF for PDF
import google.generativeai as genai
import jwt
from dotenv import load_dotenv
from flask import Flask, request, jsonify, Response, send_file
from flask_cors import CORS  # type: ignore
from flask_mongoengine import MongoEngine
from pptx import Presentation  # python-pptx for PowerPoint

# ---------------- INITIALIZATION & CONFIG ----------------
load_dotenv()

app = Flask(__name__)

CORS(app,
     origins=[
         "http://localhost:3000",
         "http://127.0.0.1:3000",
         "http://localhost:5173",
         "http://127.0.0.1:5173"
     ],
     allow_headers=["Content-Type", "x-auth-token", "Authorization"],
     supports_credentials=True,
     methods=["GET", "POST", "PATCH", "DELETE", "OPTIONS"])

app.config['MONGODB_SETTINGS'] = {
    'db': 'mindvault_db',
    'host': os.getenv('MONGO_URI')
}
app.config['JWT_SECRET'] = os.getenv('JWT_SECRET') or 'secret-dev'
db = MongoEngine(app)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
MY_VAULT_FOLDER = os.path.join(BASE_DIR, "myvault_files")
os.makedirs(MY_VAULT_FOLDER, exist_ok=True)

gemini_api_key = os.getenv("GEMINI_API_KEY")
if not gemini_api_key:
    print("⚠️ WARNING: GEMINI_API_KEY environment variable not set.")
genai.configure(api_key=gemini_api_key)

# ---------------- TOKEN DECORATOR ----------------
def token_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        auth_header = request.headers.get('Authorization')
        token = None
        if auth_header and auth_header.startswith('Bearer '):
            token = auth_header.split(' ')[1]
        else:
            token = request.headers.get('x-auth-token')

        if not token:
            return jsonify({'error': 'Token is missing!'}), 401

        try:
            data = jwt.decode(token, app.config['JWT_SECRET'], algorithms=["HS256"])
            current_user = User.objects.get(id=data['user_id'])
        except Exception as e:
            return jsonify({'error': 'Token is invalid or expired!', 'details': str(e)}), 401

        return f(current_user, *args, **kwargs)
    return decorated

# ---------------- MODELS ----------------
class User(db.Document):
    firstName = db.StringField(required=True)
    email = db.StringField(required=True, unique=True)
    password = db.StringField(required=True)
    meta = {'collection': 'users'}

class File(db.Document):
    user_id = db.ReferenceField(User, required=True)
    filename = db.StringField(required=True)
    file_id = db.StringField(required=True, unique=True)
    storage_path = db.StringField(required=True)
    file_type = db.StringField()
    mime_type = db.StringField(required=True)
    size = db.IntField()
    upload_date = db.DateTimeField(default=datetime.datetime.utcnow)
    is_deleted = db.BooleanField(default=False)
    meta = {
        'collection': 'files',
        'indexes': [
            {'fields': ('user_id', 'is_deleted')},
            {'fields': ('user_id', 'filename')},
        ]
    }

# NEW: FileChat - one chat per file per user (P1)
class FileChat(db.Document):
    user_id = db.ReferenceField(User, required=True)
    file_id = db.StringField(required=True)
    # messages: list of dicts like {role: 'user'|'assistant', text: '...', time: '...'}
    messages = db.ListField(db.DictField(), default=[])
    updated_at = db.DateTimeField(default=datetime.datetime.utcnow)

    meta = {
        'collection': 'file_chats',
        'indexes': [
            {'fields': ('user_id', 'file_id'), 'unique': True}
        ]
    }

# Other models (PlannerTask, PlannerEvent, Alert, AIPlan) unchanged from your previous file.
class PlannerTask(db.Document):
    user_id = db.ReferenceField(User, required=True)
    title = db.StringField(required=True)
    details = db.StringField()
    done = db.BooleanField(default=False)
    created_at = db.DateTimeField(default=datetime.datetime.utcnow)
    updated_at = db.DateTimeField(default=datetime.datetime.utcnow)
    meta = {'collection': 'planner_tasks'}

class PlannerEvent(db.Document):
    user_id = db.ReferenceField(User, required=True)
    title = db.StringField(required=True)
    description = db.StringField()
    deadline = db.DateTimeField(required=True)
    created_at = db.DateTimeField(default=datetime.datetime.utcnow)
    updated_at = db.DateTimeField(default=datetime.datetime.utcnow)
    meta = {'collection': 'planner_events'}

class Alert(db.Document):
    user_id = db.ReferenceField(User, required=True)
    message = db.StringField(required=True)
    related_event = db.ReferenceField(PlannerEvent, null=True)
    created_at = db.DateTimeField(default=datetime.datetime.utcnow)
    read = db.BooleanField(default=False)
    meta = {'collection': 'planner_alerts'}

class AIPlan(db.Document):
    user_id = db.ReferenceField(User, required=True)
    prompt = db.StringField()
    plan_text = db.StringField()
    created_at = db.DateTimeField(default=datetime.datetime.utcnow)
    meta = {'collection': 'ai_plans'}

# ---------------- HELPERS ----------------
def parse_iso(dt_str):
    try:
        return datetime.datetime.fromisoformat(dt_str.replace('Z', ''))
    except Exception:
        return None

# ---------------- FILE/VAULT ROUTES ----------------
@app.route('/api/vault/files', methods=['GET'])
@token_required
def get_user_files(current_user):
    try:
        search_query = request.args.get('search', '').strip()
        query_set = File.objects(user_id=current_user.id, is_deleted=False)
        if search_query:
            query_set = query_set.filter(filename__icontains=search_query)
        files = query_set.order_by('-upload_date')
        file_list = []
        for file in files:
            file_list.append({
                'id': file.file_id,
                'name': file.filename,
                'type': file.file_type,
                'mime_type': file.mime_type,
                'size': file.size,
                'date': file.upload_date.isoformat(),
            })
        return jsonify(file_list), 200
    except Exception as e:
        print(f"❌ get_user_files error: {e}")
        return jsonify({'error': 'Could not fetch files'}), 500

@app.route('/api/vault/file/<file_id>/content', methods=['GET'])
@token_required
def get_file_content(current_user, file_id):
    try:
        file_metadata = File.objects(file_id=file_id, user_id=current_user.id, is_deleted=False).first()
        if not file_metadata:
            return jsonify({"error": "File not found or access denied"}), 404
        target_file_path = file_metadata.storage_path
        if not os.path.exists(target_file_path):
            return jsonify({"error": "File content not on server"}), 404
        return send_file(
            target_file_path,
            mimetype=file_metadata.mime_type,
            as_attachment=False
        )
    except Exception as e:
        print("❌ get_file_content error:", e)
        return jsonify({"error": str(e)}), 500

# ---------------- CHAT ENDPOINTS ----------------
@app.route('/api/chat/<file_id>', methods=['GET'])
@token_required
def get_chat(current_user, file_id):
    """Return saved chat for this file (one chat per file)."""
    try:
        chat_doc = FileChat.objects(user_id=current_user, file_id=file_id).first()
        if not chat_doc:
            return jsonify({"messages": []}), 200
        return jsonify({"messages": chat_doc.messages or []}), 200
    except Exception as e:
        print("❌ get_chat error:", e)
        return jsonify({"error": "Could not fetch chat"}), 500

@app.route('/api/chat/<file_id>/save', methods=['POST'])
@token_required
def save_chat(current_user, file_id):
    """Save (upsert) the chat messages for the file."""
    try:
        body = request.get_json(force=True)
        messages = body.get('messages', [])
        if not isinstance(messages, list):
            return jsonify({"error": "Invalid messages format"}), 400

        # Upsert: if exists update, else create
        chat_doc = FileChat.objects(user_id=current_user, file_id=file_id).first()
        if chat_doc:
            chat_doc.messages = messages
            chat_doc.updated_at = datetime.datetime.utcnow()
            chat_doc.save()
        else:
            FileChat(user_id=current_user, file_id=file_id, messages=messages).save()

        return jsonify({"ok": True}), 200
    except Exception as e:
        print("❌ save_chat error:", e)
        return jsonify({"error": str(e)}), 500

@app.route('/api/chat/<file_id>/ask', methods=['POST'])
@token_required
def ask_chat(current_user, file_id):
    """
    Ask a question about the file. We will:
    - load saved chat for context (if any)
    - load file text (first N chars) as additional context
    - call Gemini to answer
    - return answer (but do NOT auto-save)
    """
    try:
        body = request.get_json(force=True)
        question = (body.get('question') or '').strip()
        if not question:
            return jsonify({"error": "Missing question"}), 400

        # 1) Load saved chat for context
        chat_doc = FileChat.objects(user_id=current_user, file_id=file_id).first()
        saved_messages = chat_doc.messages if chat_doc else []

        # 2) load file content (first chunk) to provide context
        file_meta = File.objects(file_id=file_id, user_id=current_user.id, is_deleted=False).first()
        file_text = ""
        if file_meta and os.path.exists(file_meta.storage_path):
            path = file_meta.storage_path
            try:
                if path.endswith(".pdf"):
                    with fitz.open(path) as doc:
                        # read first few pages (avoid huge prompts)
                        pages_to_read = min(5, doc.page_count)
                        file_text = "\n".join([doc.load_page(i).get_text() for i in range(pages_to_read)])
                elif path.endswith((".pptx", ".ppt")):
                    prs = Presentation(path)
                    texts = []
                    for slide in prs.slides[:10]:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                texts.append(shape.text)
                    file_text = "\n".join(texts)
                elif path.endswith(".txt"):
                    with open(path, "r", encoding="utf-8") as f:
                        file_text = f.read(4000)
            except Exception as e:
                print("⚠️ file_text read error:", e)
                file_text = ""

        # 3) Build prompt for Gemini (simple concatenation)
        # Use saved_messages to keep conversational context
        conversation_context = ""
        for m in saved_messages[-10:]:  # keep last 10 messages
            role = m.get('role', 'user')
            text = m.get('text', '')
            conversation_context += f"{role.upper()}: {text}\n"

        prompt = (
            "You are an assistant that answers questions about the uploaded file.\n\n"
            f"File content (first chunk):\n{file_text[:4000]}\n\n"
            f"Conversation so far:\n{conversation_context}\n"
            f"User: {question}\n\nAnswer concisely and helpfully."
        )

        # 4) Call Gemini model
        try:
            model = genai.GenerativeModel("gemini-2.5")
            resp = model.generate_content(prompt)
            answer = getattr(resp, 'text', None)
            if not answer and hasattr(resp, 'candidates'):
                answer = resp.candidates[0].content.parts[0].text
            if not answer:
                answer = "No answer generated."
        except Exception as e:
            print("⚠️ Gemini ask error:", e)
            # fallback: simple canned response
            answer = "Sorry — I couldn't reach the AI model right now."

        # 5) return answer (do not save automatically)
        return jsonify({"answer": answer}), 200

    except Exception as e:
        print("❌ ask_chat error:", e)
        return jsonify({"error": str(e)}), 500

# ---------------- PLANNER ROUTES (UNCHANGED: generate-plan, tasks etc.) ----------------
@app.route('/api/planner/generate-plan', methods=['POST'])
@token_required
def generate_plan(current_user):
    try:
        try:
            body = request.get_json(force=True)
            if isinstance(body, str):
                body = json.loads(body)
        except Exception:
            body = {}

        goals = (body.get('goals') or '').strip()
        subjects = (body.get('subjects') or '').strip()
        timeframe = (body.get('timeframe') or '').strip()

        if not (goals or subjects or timeframe):
            return jsonify({"error": "Provide at least one of goals/subjects/timeframe"}), 400

        prompt = (
            f"You are an expert study planner. Create a clear, actionable study plan.\n\n"
            f"Goals: {goals}\nSubjects: {subjects}\nTimeframe: {timeframe}\n\n"
            "Return the plan in markdown format with day-wise steps."
        )

        plan_text = None
        try:
            model = genai.GenerativeModel('gemini-2.5')
            resp = model.generate_content(prompt)
            plan_text = getattr(resp, 'text', None)
            if not plan_text and hasattr(resp, 'candidates'):
                plan_text = resp.candidates[0].content.parts[0].text
        except Exception as inner_e:
            print("⚠️ Gemini error:", inner_e)
            plan_text = None

        if not plan_text:
            plan_text = (
                f"### Study Plan (auto-created)\n\n**Focus:** {subjects or goals}\n\n"
                "- Day 1: Read core concepts\n"
                "- Day 2: Practice problems\n"
                "- Day 3: Revise and summarize\n\n"
                "_This is an auto-generated fallback plan._"
            )

        AIPlan(user_id=current_user, prompt=prompt, plan_text=plan_text).save()
        return jsonify({"plan": plan_text}), 200

    except Exception as e:
        print("❌ generate_plan error:", e)
        return jsonify({"error": str(e)}), 500

# ---------------- TASKS / UPDATES (UNCHANGED) ----------------
@app.route('/api/planner/tasks', methods=['POST'])
@token_required
def create_task(current_user):
    body = request.get_json() or {}
    title = (body.get('title') or '').strip()
    details = body.get('details', '').strip()
    if not title:
        return jsonify({"error": "Missing title"}), 400
    task = PlannerTask(user_id=current_user, title=title, details=details)
    task.save()
    return jsonify({"taskId": str(task.id), "title": task.title}), 201

@app.route('/api/planner/tasks', methods=['GET'])
@token_required
def get_tasks(current_user):
    tasks = PlannerTask.objects(user_id=current_user).order_by('-created_at')
    data = [{"id": str(t.id), "title": t.title, "details": t.details, "done": t.done} for t in tasks]
    return jsonify({"tasks": data})

@app.route('/api/planner/tasks/<task_id>', methods=['PATCH'])
@token_required
def update_task(current_user, task_id):
    body = request.get_json() or {}
    task = PlannerTask.objects(id=task_id, user_id=current_user).first()
    if not task:
        return jsonify({"error": "Not found"}), 404
    for key in ['title', 'details']:
        if key in body:
            setattr(task, key, body[key])
    if 'done' in body:
        task.done = bool(body['done'])
    task.updated_at = datetime.datetime.utcnow()
    task.save()
    return jsonify({"ok": True})

@app.route('/api/planner/events', methods=['POST'])
@token_required
def create_event(current_user):
    body = request.get_json() or {}
    title = (body.get('title') or '').strip()
    deadline = body.get('deadline')
    description = body.get('description', '').strip()
    if not title or not deadline:
        return jsonify({"error": "Missing fields"}), 400
    dt = parse_iso(deadline)
    if not dt:
        return jsonify({"error": "Invalid date format"}), 400
    ev = PlannerEvent(user_id=current_user, title=title, description=description, deadline=dt)
    ev.save()
    return jsonify({"eventId": str(ev.id)}), 201

@app.route('/api/planner/events', methods=['GET'])
@token_required
def get_events(current_user):
    q = PlannerEvent.objects(user_id=current_user)
    events = [{"id": str(e.id), "title": e.title, "description": e.description, "deadline": e.deadline.isoformat()} for e in q]
    return jsonify({"events": events})

@app.route('/api/planner/upcoming-deadlines', methods=['GET'])
@token_required
def get_upcoming_deadlines(current_user):
    now = datetime.datetime.utcnow()
    evs = PlannerEvent.objects(user_id=current_user, deadline__gt=now).order_by('deadline')
    out = [{"id": str(e.id), "title": e.title, "deadline": e.deadline.isoformat()} for e in evs]
    return jsonify({"deadlines": out})

@app.route('/api/planner/alerts', methods=['GET'])
@token_required
def get_alerts(current_user):
    now = datetime.datetime.utcnow()
    in_one_hour = now + datetime.timedelta(hours=1)

    expired = PlannerEvent.objects(user_id=current_user, deadline__lte=now, description__ne="Reminder")
    soon = PlannerEvent.objects(user_id=current_user, deadline__gt=now, deadline__lte=in_one_hour, description__ne="Reminder")

    alerts = [
        {
            "id": str(e.id),
            "type": "expired",
            "message": f"Deadline expired: {e.title}",
            "deadline": e.deadline.isoformat(),
        }
        for e in expired
    ]
    alerts += [
        {
            "id": str(e.id),
            "type": "upcoming",
            "message": f"Upcoming soon: {e.title}",
            "deadline": e.deadline.isoformat(),
        }
        for e in soon
    ]

    reminders = PlannerEvent.objects(user_id=current_user, description="Reminder")
    alerts += [
        {
            "id": str(r.id),
            "type": "reminder",
            "message": f"Reminder: {r.title}",
            "deadline": r.deadline.isoformat() if r.deadline else None,
        }
        for r in reminders
    ]

    return jsonify({"alerts": alerts})

# ---------------- BACKGROUND CHECKER ----------------
def planner_background_checker(interval=30):
    def run():
        with app.app_context():
            while True:
                try:
                    now = datetime.datetime.utcnow()
                    expired_events = PlannerEvent.objects(deadline__lte=now)
                    for ev in expired_events:
                        if not Alert.objects(related_event=ev).first():
                            Alert(user_id=ev.user_id, message=f"Deadline expired: {ev.title}", related_event=ev).save()
                except Exception as e:
                    print("Planner background checker error:", e)
                time.sleep(interval)
    threading.Thread(target=run, daemon=True).start()

planner_background_checker(30)

# ---------------- AUTH (unchanged) ----------------
@app.route('/api/auth/register', methods=['POST'])
def register_user():
    body = request.get_json()
    if not body or not body.get('email') or not body.get('password'):
        return jsonify({"error": "Missing required fields"}), 400
    if User.objects(email=body.get('email')).first():
        return jsonify({"error": "User exists"}), 409
    hashed = bcrypt.hashpw(body['password'].encode('utf-8'), bcrypt.gensalt())
    user = User(firstName=body.get('firstName'), email=body['email'], password=hashed.decode('utf-8')).save()
    return jsonify({"message": f"User '{user.email}' registered successfully"}), 201

@app.route('/api/auth/login', methods=['POST'])
def login_user():
    body = request.get_json()
    user = User.objects(email=body.get('email')).first()
    if not user or not bcrypt.checkpw(body['password'].encode('utf-8'), user.password.encode('utf-8')):
        return jsonify({"error": "Invalid credentials"}), 401
    token = jwt.encode({'user_id': str(user.id), 'exp': datetime.datetime.utcnow() + datetime.timedelta(hours=5)}, app.config['JWT_SECRET'], algorithm="HS256")
    return jsonify({"token": token})

@app.route('/api/auth/me', methods=['GET'])
@token_required
def get_current_user(current_user):
    return jsonify(id=str(current_user.id), firstName=current_user.firstName, email=current_user.email)

# ---------------- FILE UPLOAD & AI ROUTES (upload, summarize, mcqs) ----------------
@app.route('/api/upload', methods=['POST'])
@token_required
def upload_file(current_user):
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    file_type = request.form.get('type', 'Unknown')
    if not file.filename:
        return jsonify({"error": "Empty filename"}), 400

    original_filename = file.filename
    file_id = str(uuid.uuid4())
    ext = os.path.splitext(original_filename)[1]

    save_path = os.path.join(MY_VAULT_FOLDER, f"{file_id}{ext}")
    file.save(save_path)
    file_size = os.path.getsize(save_path)

    new_file = File(
        user_id=current_user,
        filename=original_filename,
        file_id=file_id,
        storage_path=save_path,
        file_type=file_type,
        mime_type=file.mimetype,
        size=file_size
    )
    new_file.save()

    response_file = {
        'id': new_file.file_id,
        'name': new_file.filename,
        'type': new_file.file_type,
        'mime_type': new_file.mime_type,
        'size': new_file.size,
        'date': new_file.upload_date.isoformat(),
    }

    return jsonify({"fileId": file_id, "file": response_file}), 200

@app.route('/api/summarize/<file_id>', methods=['GET'])
@token_required
def summarize_file(current_user, file_id):
    try:
        file_metadata = File.objects(file_id=file_id, user_id=current_user.id).first()
        if not file_metadata:
            return jsonify({"error": "File not found in vault"}), 404

        target_file = file_metadata.storage_path
        if not os.path.exists(target_file):
            return jsonify({"error": "File content not on server"}), 404

        text_content = ""
        if target_file.endswith(".pdf"):
            with fitz.open(target_file) as doc:
                text_content = "\n".join([page.get_text() for page in doc])
        elif target_file.endswith((".pptx", ".ppt")):
            prs = Presentation(target_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
        elif target_file.endswith(".txt"):
            with open(target_file, "r", encoding="utf-8") as f:
                text_content = f.read()
        else:
            return jsonify({"error": "Unsupported file type"}), 400

        model = genai.GenerativeModel("gemini-2.0-flash")
        prompt = f"Summarize this text concisely:\n\n{text_content[:5000]}"
        response = model.generate_content(prompt)
        summary = response.text.strip() if hasattr(response, "text") else "No summary generated."

        return jsonify({"summary": summary}), 200

    except Exception as e:
        print("❌ summarize_file error:", e)
        return jsonify({"error": str(e)}), 500

@app.route('/api/mcqs/<file_id>', methods=['GET'])
@token_required
def generate_mcqs(current_user, file_id):
    try:
        file_metadata = File.objects(file_id=file_id, user_id=current_user.id).first()
        if not file_metadata:
            return jsonify({"error": "File not found in vault"}), 404

        target_file = file_metadata.storage_path
        if not os.path.exists(target_file):
            return jsonify({"error": "File content not on server"}), 404

        text_content = ""
        if target_file.endswith(".pdf"):
            with fitz.open(target_file) as doc:
                text_content = "\n".join([page.get_text() for page in doc])
        elif target_file.endswith((".pptx", ".ppt")):
            prs = Presentation(target_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
        elif target_file.endswith(".txt"):
            with open(target_file, "r", encoding="utf-8") as f:
                text_content = f.read()

        model = genai.GenerativeModel("gemini-2.0-flash")
        prompt = f"Generate 5 MCQs from this content with options and correct answers in JSON format:\n{text_content[:5000]}"
        response = model.generate_content(prompt)

        mcqs_json = []
        try:
            cleaned_text = response.text.strip().replace('```json', '').replace('```', '')
            mcqs_json = json.loads(cleaned_text)
        except Exception:
            mcqs_json = []

        return jsonify({"mcqs": mcqs_json}), 200

    except Exception as e:
        print("❌ generate_mcqs error:", e)
        return jsonify({"error": str(e)}), 500

@app.route('/api/vault/file/<file_id>/delete', methods=['DELETE'])
@token_required
def delete_file_permanently(current_user, file_id):
    try:
        file = File.objects(file_id=file_id, user_id=current_user.id).first()
        if not file:
            return jsonify({"error": "File not found"}), 404

        # Remove file from storage
        if os.path.exists(file.storage_path):
            os.remove(file.storage_path)

        # Remove from database
        file.delete()

        # Also delete chat for this file (optional but correct)
        FileChat.objects(file_id=file_id, user_id=current_user.id).delete()

        return jsonify({"success": True}), 200

    except Exception as e:
        print("❌ Permanent delete error:", e)
        return jsonify({"error": str(e)}), 500

# ---------------- RUN APP ----------------
if __name__ == '__main__':
    if not any(t.name == 'PlannerChecker' for t in threading.enumerate()):
        checker_thread = threading.Thread(target=planner_background_checker, args=(30,), daemon=True, name='PlannerChecker')
        checker_thread.start()

    app.run(host='0.0.0.0', port=5000, debug=True)